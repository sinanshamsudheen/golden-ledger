"""
Document text extraction module.

Supports:
  - PDF  (.pdf)   via pdfminer.six
  - PPTX (.pptx)  via python-pptx
  - DOCX (.docx)  via python-docx
"""

import io
import logging

logger = logging.getLogger(__name__)


class PasswordProtectedError(Exception):
    """Raised when a document is password-protected and cannot be parsed."""


# OLE2 magic bytes — encrypted Office files (docx/pptx/xlsx) are CFBF containers
_OLE_MAGIC = b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"


def _is_ole_encrypted(content: bytes) -> bool:
    return content[:8] == _OLE_MAGIC


def extract_text(content: bytes, file_name: str) -> str:
    """
    Dispatch to the correct extractor based on file extension.

    Args:
        content:   Raw bytes of the document file.
        file_name: Original file name (used to determine format).

    Returns:
        Extracted plain-text string.

    Raises:
        ValueError: If the file format is not supported.
    """
    lower = file_name.lower()
    if lower.endswith(".pdf"):
        return _extract_pdf(content)
    elif lower.endswith(".pptx"):
        return _extract_pptx(content)
    elif lower.endswith(".docx"):
        return _extract_docx(content)
    elif lower.endswith(".txt"):
        return content.decode("utf-8", errors="ignore")
    else:
        raise ValueError(f"Unsupported file format: '{file_name}'")


# ── PDF ───────────────────────────────────────────────────────────────────────

def _extract_pdf(content: bytes) -> str:
    from pdfminer.high_level import extract_text_to_fp
    from pdfminer.layout import LAParams

    try:
        from pdfminer.pdfdocument import PDFPasswordIncorrect as _PDFPasswordIncorrect
    except ImportError:
        _PDFPasswordIncorrect = None  # type: ignore[assignment]

    output = io.StringIO()
    try:
        extract_text_to_fp(
            io.BytesIO(content),
            output,
            laparams=LAParams(),
            output_type="text",
            codec="utf-8",
        )
    except Exception as exc:
        # Detect by class name for forward-compat with pdfminer version differences
        if (
            (_PDFPasswordIncorrect and isinstance(exc, _PDFPasswordIncorrect))
            or type(exc).__name__ == "PDFPasswordIncorrect"
        ):
            raise PasswordProtectedError("PDF is password-protected") from exc
        raise
    return output.getvalue().strip()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _extract_pptx(content: bytes) -> str:
    if _is_ole_encrypted(content):
        raise PasswordProtectedError("PPTX is password-protected")
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
    return "\n".join(lines)


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _extract_docx(content: bytes) -> str:
    if _is_ole_encrypted(content):
        raise PasswordProtectedError("DOCX is password-protected")
    import docx

    doc = docx.Document(io.BytesIO(content))
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
